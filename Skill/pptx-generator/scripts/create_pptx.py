from pptx import Presentation
from pptx.util import Inches, Pt, Cm
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE
import sys
import json
import argparse
import os

# ── Color Palette ─────────────────────────────────────────────────────────────
NAVY = RGBColor(0x1F, 0x38, 0x64)      # Primary brand color
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
BLACK = RGBColor(0x00, 0x00, 0x00)
LIGHT_GRAY = RGBColor(0xF2, 0xF2, 0xF2)

# ── Font Config ────────────────────────────────────────────────────────────────
FONT_TITLE = "Calibri Light"
FONT_BODY = "Calibri"


def _set_slide_background(slide, color: RGBColor):
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = color


def _add_header_bar(slide, prs):
    width = prs.slide_width
    bar = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(0), Inches(0),
        width, Cm(1.2)
    )
    bar.fill.solid()
    bar.fill.fore_color.rgb = NAVY
    bar.line.fill.background()


def _add_footer(slide, text="", page_num=None):
    txBox = slide.shapes.add_textbox(
        Inches(0.3), Inches(6.9),
        Inches(9.4), Inches(0.4)
    )
    tf = txBox.text_frame
    tf.paragraphs[0].text = text
    tf.paragraphs[0].font.size = Pt(10)
    tf.paragraphs[0].font.color.rgb = RGBColor(0x80, 0x80, 0x80)


def _add_title_slide(slide, config, prs):
    _set_slide_background(slide, WHITE)
    _add_header_bar(slide, prs)

    # Title
    txBox = slide.shapes.add_textbox(Inches(0.8), Inches(2.2), Inches(8.4), Inches(1.5))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = config.get("title", "")
    p.font.name = FONT_TITLE
    p.font.size = Pt(44)
    p.font.bold = True
    p.font.color.rgb = NAVY
    p.alignment = PP_ALIGN.CENTER

    # Subtitle
    sub = config.get("subtitle", "")
    if sub:
        txBox2 = slide.shapes.add_textbox(Inches(0.8), Inches(3.8), Inches(8.4), Inches(1.0))
        tf2 = txBox2.text_frame
        tf2.word_wrap = True
        p2 = tf2.paragraphs[0]
        p2.text = sub
        p2.font.name = FONT_BODY
        p2.font.size = Pt(24)
        p2.font.color.rgb = RGBColor(0x44, 0x44, 0x44)
        p2.alignment = PP_ALIGN.CENTER


def _add_bullet_slide(slide, config, prs):
    _set_slide_background(slide, WHITE)
    _add_header_bar(slide, prs)

    # Title
    title_text = config.get("title", "")
    txTitle = slide.shapes.add_textbox(Inches(0.8), Inches(1.5), Inches(8.4), Inches(1.0))
    tf_t = txTitle.text_frame
    tf_t.word_wrap = True
    p_t = tf_t.paragraphs[0]
    p_t.text = title_text
    p_t.font.name = FONT_TITLE
    p_t.font.size = Pt(32)
    p_t.font.bold = True
    p_t.font.color.rgb = NAVY

    # Bullets
    bullets = config.get("bullets", [])
    txBox = slide.shapes.add_textbox(Inches(0.8), Inches(2.7), Inches(8.4), Inches(3.5))
    tf = txBox.text_frame
    tf.word_wrap = True

    for i, bullet in enumerate(bullets):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = f"• {bullet}"
        p.font.name = FONT_BODY
        p.font.size = Pt(22)
        p.font.color.rgb = BLACK
        p.space_before = Pt(8)
        p.level = 0


def _add_two_col_slide(slide, config, prs):
    _set_slide_background(slide, WHITE)
    _add_header_bar(slide, prs)

    # Title
    txTitle = slide.shapes.add_textbox(Inches(0.8), Inches(1.5), Inches(8.4), Inches(0.8))
    tf_t = txTitle.text_frame
    p_t = tf_t.paragraphs[0]
    p_t.text = config.get("title", "")
    p_t.font.name = FONT_TITLE
    p_t.font.size = Pt(32)
    p_t.font.bold = True
    p_t.font.color.rgb = NAVY

    # Left column
    txLeft = slide.shapes.add_textbox(Inches(0.5), Inches(2.5), Inches(4.3), Inches(4.0))
    tf_l = txLeft.text_frame
    tf_l.word_wrap = True
    p_l = tf_l.paragraphs[0]
    p_l.text = config.get("col_left", "")
    p_l.font.name = FONT_BODY
    p_l.font.size = Pt(20)
    p_l.font.color.rgb = BLACK

    # Right column
    txRight = slide.shapes.add_textbox(Inches(5.2), Inches(2.5), Inches(4.3), Inches(4.0))
    tf_r = txRight.text_frame
    tf_r.word_wrap = True
    p_r = tf_r.paragraphs[0]
    p_r.text = config.get("col_right", "")
    p_r.font.name = FONT_BODY
    p_r.font.size = Pt(20)
    p_r.font.color.rgb = BLACK

    # Vertical divider
    line = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(5.0), Inches(2.5),
        Pt(1), Inches(3.5)
    )
    line.fill.solid()
    line.fill.fore_color.rgb = NAVY
    line.line.fill.background()


def _add_blank_slide(slide, config, prs):
    _set_slide_background(slide, WHITE)
    _add_header_bar(slide, prs)

    if config.get("title"):
        txBox = slide.shapes.add_textbox(Inches(0.8), Inches(2.0), Inches(8.4), Inches(3.0))
        tf = txBox.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = config["title"]
        p.font.name = FONT_BODY
        p.font.size = Pt(28)
        p.font.color.rgb = BLACK
        p.alignment = PP_ALIGN.CENTER


# ── Main ───────────────────────────────────────────────────────────────────────
def create_presentation(slides_config: list, output_path: str):
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)

    slide_type_map = {
        "title": _add_title_slide,
        "bullet": _add_bullet_slide,
        "two_col": _add_two_col_slide,
        "blank": _add_blank_slide,
    }

    for config in slides_config:
        slide_type = config.get("type", "blank")
        blank_layout = prs.slide_layouts[6]
        slide = prs.slides.add_slide(blank_layout)

        handler = slide_type_map.get(slide_type, _add_blank_slide)
        handler(slide, config, prs)

    prs.save(output_path)
    return output_path


def main():
    parser = argparse.ArgumentParser(description="Tạo file PowerPoint từ JSON config")
    parser.add_argument("--output", "-o", required=True, help="Đường dẫn file .pptx đầu ra")
    parser.add_argument("--slides", "-s", required=True, help="JSON string chứa mảng slide config")
    args = parser.parse_args()

    try:
        slides_config = json.loads(args.slides)
    except json.JSONDecodeError:
        slides_config = json.loads(args.slides.replace("'", '"'))
    except Exception:
        import ast
        try:
            slides_config = ast.literal_eval(args.slides)
        except Exception as e:
            print(f"Lỗi parse JSON: {e}", file=sys.stderr)
            sys.exit(1)

    if not isinstance(slides_config, list):
        slides_config = [slides_config]

    output_path = os.path.abspath(args.output)
    os.makedirs(os.path.dirname(output_path) or ".", exist_ok=True)

    create_presentation(slides_config, output_path)
    print(f"[pptx-generator] Created: {output_path}")


if __name__ == "__main__":
    main()